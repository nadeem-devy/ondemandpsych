"""
Document text extractors for PDF, PPT/PPTX, DOCX, and Video/Audio files.
Handles the 2K+ PDFs, 5K+ PPTs, and 1.5K+ videos at scale.
"""

import os
import tempfile
import subprocess
import logging
from pathlib import Path

import pdfplumber
from pptx import Presentation
from docx import Document

logger = logging.getLogger(__name__)


def extract_pdf(file_path: str) -> tuple[str, int]:
    """Extract text from PDF. Returns (text, page_count)."""
    text_parts = []
    page_count = 0

    with pdfplumber.open(file_path) as pdf:
        page_count = len(pdf.pages)
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text_parts.append(page_text)

            # Also extract tables as text
            tables = page.extract_tables()
            for table in tables:
                for row in table:
                    cells = [str(cell) if cell else "" for cell in row]
                    text_parts.append(" | ".join(cells))

    return "\n\n".join(text_parts), page_count


def extract_pptx(file_path: str) -> tuple[str, int]:
    """Extract text from PowerPoint. Returns (text, slide_count)."""
    prs = Presentation(file_path)
    text_parts = []
    slide_count = len(prs.slides)

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = [f"[Slide {slide_num}]"]

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_texts.append(text)

            # Extract from tables
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    slide_texts.append(" | ".join(cells))

            # Extract from grouped shapes
            if hasattr(shape, "shapes"):
                for sub_shape in shape.shapes:
                    if sub_shape.has_text_frame:
                        for paragraph in sub_shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text:
                                slide_texts.append(text)

        # Extract notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_texts.append(f"[Notes] {notes}")

        text_parts.append("\n".join(slide_texts))

    return "\n\n".join(text_parts), slide_count


def extract_docx(file_path: str) -> tuple[str, int]:
    """Extract text from Word document. Returns (text, page_estimate)."""
    doc = Document(file_path)
    text_parts = []

    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if text:
            text_parts.append(text)

    # Extract from tables
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            text_parts.append(" | ".join(cells))

    full_text = "\n\n".join(text_parts)
    page_estimate = max(1, len(full_text) // 3000)  # rough estimate
    return full_text, page_estimate


def extract_video_audio(file_path: str, whisper_model: str = "base") -> tuple[str, int]:
    """
    Extract transcript from video/audio using OpenAI Whisper.
    Returns (transcript, duration_seconds).
    """
    import whisper

    # Extract audio from video if needed
    audio_path = file_path
    temp_audio = None

    video_extensions = {".mp4", ".avi", ".mov", ".mkv", ".webm", ".flv", ".wmv"}
    ext = Path(file_path).suffix.lower()

    if ext in video_extensions:
        temp_audio = tempfile.NamedTemporaryFile(suffix=".wav", delete=False)
        temp_audio.close()
        audio_path = temp_audio.name

        try:
            subprocess.run(
                [
                    "ffmpeg", "-i", file_path,
                    "-vn", "-acodec", "pcm_s16le",
                    "-ar", "16000", "-ac", "1",
                    "-y", audio_path,
                ],
                capture_output=True,
                check=True,
                timeout=600,  # 10 min max
            )
        except FileNotFoundError:
            raise RuntimeError(
                "ffmpeg not found. Install ffmpeg to process video files: "
                "brew install ffmpeg (macOS) or apt-get install ffmpeg (Linux)"
            )
        except subprocess.TimeoutExpired:
            raise RuntimeError("Video audio extraction timed out (>10 minutes)")

    try:
        model = whisper.load_model(whisper_model)
        result = model.transcribe(audio_path)

        transcript = result["text"]
        # Get duration from segments
        duration = 0
        if result.get("segments"):
            duration = int(result["segments"][-1].get("end", 0))

        return transcript, duration
    finally:
        if temp_audio and os.path.exists(temp_audio.name):
            os.unlink(temp_audio.name)


def extract_text(file_path: str, file_type: str, whisper_model: str = "base") -> tuple[str, int]:
    """
    Route to the correct extractor based on file type.
    Returns (extracted_text, page_or_slide_count).
    """
    file_type = file_type.lower().strip(".")

    extractors = {
        "pdf": extract_pdf,
        "pptx": extract_pptx,
        "ppt": extract_pptx,  # python-pptx handles .ppt conversion
        "docx": extract_docx,
        "doc": extract_docx,
        "txt": lambda fp: (Path(fp).read_text(encoding="utf-8", errors="replace"), 1),
        "csv": lambda fp: (Path(fp).read_text(encoding="utf-8", errors="replace"), 1),
    }

    video_types = {"mp4", "avi", "mov", "mkv", "webm", "flv", "wmv", "mp3", "wav", "m4a", "ogg"}

    if file_type in extractors:
        return extractors[file_type](file_path)
    elif file_type in video_types:
        return extract_video_audio(file_path, whisper_model)
    else:
        raise ValueError(f"Unsupported file type: {file_type}")
